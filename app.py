import os
import base64
import requests
from openai import AzureOpenAI
from flask import Flask, redirect, render_template, request, send_from_directory, url_for, abort
from docx import Document
from unidecode import unidecode
import fitz  # PyMuPDF
from pptx import Presentation

# OpenAI API configuration
client = AzureOpenAI(
    azure_endpoint="insert_your_azure_key_endpoint_here",
    api_key='insert_your_azure_key_here',
    api_version="insert_your_api_version_here" #"2024-02-01"
)

# Function to extract text and images from DOCX files
def extractContentDocs(file):
    document = Document(file)
    content = []
    for para in document.paragraphs:
        content.append(para.text)  # Adds the text of each paragraph to the content
    
    for rel in document.part.rels.values():
        if "image" in rel.target_ref:
            image = rel.target_part.blob  # Adds images to the content
            content.append(image)
    
    return content

# Function to extract text and images from PDF files
def extractContentPdf(file):
    document = fitz.open(stream=file.read(), filetype="pdf")
    content = []
    for page in document:
        text = page.get_text()  # Adds the text of each page to the content
        content.append(text)
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = document.extract_image(xref)
            image_bytes = base_image["image"]  # Adds images to the content
            content.append(image_bytes)
    return content

# Function to extract text and images from PPT/PPTX files
def extractContentPpt(file):
    prs = Presentation(file)
    content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text  # Adds the text of each slide to the content
                content.append(text)
            elif hasattr(shape, "image"):
                image = shape.image.blob  # Adds images to the content
                content.append(image)
    return content

# Function to upload image to ImgBB and get the public URL
def upload_image_to_imgbb(image_bytes):
    imgbb_api_key = 'insert_your_imgbb_key_here'
    imgbb_url = 'https://api.imgbb.com/1/upload'
    image_base64 = base64.b64encode(image_bytes).decode('utf-8')
    payload = {
        'key': imgbb_api_key,
        'image': image_base64
    }
    imgbb_response = requests.post(imgbb_url, data=payload)
    if imgbb_response.status_code == 200:
        imgbb_data = imgbb_response.json()
        public_image_url = imgbb_data['data']['url']  # Public URL of the image
        delete_url = imgbb_data['data']['delete_url']  # URL to delete the image
        return public_image_url, delete_url
    return None, None

# Function to describe image using OpenAI
def describe_image(image_url):
    response = client.chat.completions.create(
        model="insert_your_openai_model_implementation_here",
        messages=[
            { "role": "system", "content": "You are an excellent assistant for describing images." },
            { "role": "user", "content": [  
                { 
                    "type": "text", 
                    "text": "Describe this image:" 
                },
                { 
                    "type": "image_url",
                    "image_url": {
                        "url": image_url
                    }
                }
            ] } 
        ],
    )
    return response.choices[0].message.content  # Returns the image description

# Function to extract and convert text from TXT files
def extractTextTxts(file):
    try:
        text = file.read().decode("utf-8")
        return unidecode(text)  # Removes accents from the text
    except FileNotFoundError:
        print(f"Error: File not found.")
        return None

# Initialize the Flask application
app = Flask(__name__)

# Route for the main page
@app.route('/')
def index():
    print('Request for index page received')
    return render_template('index.html')

# Route for the favicon
@app.route('/favicon.ico')
def favicon():
    return send_from_directory(os.path.join(app.root_path, 'static'), 'favicon.ico', mimetype='image/vnd.microsoft.icon')

# Route to process the form submission result
@app.route('/result', methods=['POST'])
def result():
    try:
        file = request.files.get('file')
        if file and file.filename != '':
            print('Request for result page received with file')
            ext = os.path.splitext(file.filename)[1].lower()
            content = []
            image_base64 = None

            # Extract content based on file extension
            if ext == '.docx':
                content = extractContentDocs(file)
            elif ext == '.pdf':
                content = extractContentPdf(file)
            elif ext == '.txt':
                content.append(extractTextTxts(file))
            elif ext == '.pptx':
                content = extractContentPpt(file)
            elif ext in ['.png', '.jpeg', '.jpg']:
                image_bytes = file.read()
                image_base64 = base64.b64encode(image_bytes).decode('utf-8') # Encode the image in base64
                public_image_url, delete_url = upload_image_to_imgbb(image_bytes)
                if public_image_url:
                    description = describe_image(public_image_url)
                    content.append(description)
                    requests.get(delete_url) 

            final_content = ""
            # Concatenate the extracted content into a final string
            for item in content:
                if isinstance(item, str):
                    final_content += item + "\n"
                elif isinstance(item, bytes):
                    public_image_url, delete_url = upload_image_to_imgbb(item)
                    if public_image_url:
                        description = describe_image(public_image_url)
                        final_content += description + "\n"
                        requests.get(delete_url) 

            # Extract the prompt from the "prompts.txt" file for OpenAI
            prompt = extractTextTxts(open('prompts.txt', 'rb'))
            response = client.chat.completions.create(
                model="insert_your_openai_model_implementation_here",
                messages=[
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": final_content}
                ],
                max_tokens=1000 # Token limit for the response
            )
            content = response.choices[0].message.content
            content = content.replace('**', '') # Remove bold markers
            
            with open("output.txt", 'w', encoding='utf-8') as file:
                file.write(content)
            return render_template('result.html', txt=content, image_base64=image_base64)
            
        print('Request for result page received with no file or unsupported file type -- redirecting')
        return redirect(url_for('index'))

    except Exception as e:
        # Log the exception
        print(f"Error: {e}")
        # Render the result.html page with the error message
        return render_template('result.html', txt="Oops... we had a problem and could not analyze the document. Please try again later", image_base64=None)

# Error handler for unhandled exceptions
@app.errorhandler(Exception)
def handle_exception(e):
    # You can add logging here if you want
    print(f"Error: {e}")
    # Render the result.html page with the error message
    return render_template('result.html', txt="Oops... we had a problem and could not analyze the document. Please try again later", image_base64=None)

if __name__ == '__main__':
    app.run()